import os
import difflib
import openai
import pandas as pd
import tkinter as tk
from tkinter import filedialog, messagebox, ttk, StringVar, Listbox, END, Entry
from pptx import Presentation
from openpyxl import load_workbook
from datetime import datetime
from dotenv import load_dotenv

# .env 파일의 환경 변수를 로드합니다.
load_dotenv()

# API 키를 환경 변수에서 가져옵니다.
API_KEY = os.getenv("FLASK_API_KEY")

# OpenAI API 키 설정
openai.api_key = API_KEY

class AutocompleteEntry(Entry):
    def __init__(self, autocomplete_list, *args, **kwargs):
        super().__init__(*args, **kwargs)
        self.autocomplete_list = autocomplete_list
        self.var = StringVar()
        self.config(textvariable=self.var)
        self.var.trace('w', self.changed)
        self.bind("<Right>", self.selection)
        self.bind("<Double-Button-1>", self.selection)
        self.bind("<Return>", self.selection)
        self.bind("<Up>", self.move_up)
        self.bind("<Down>", self.move_down)
        self.lb_up = False

    def changed(self, name, index, mode):
        if self.var.get() == '':
            if self.lb_up:
                self.lb.destroy()
                self.lb_up = False
        else:
            words = self.comparison()
            if words:
                if not self.lb_up:
                    self.lb = Listbox(self.master)
                    self.lb.bind("<Double-Button-1>", self.selection)
                    self.lb.bind("<Right>", self.selection)
                    self.lb.bind("<Return>", self.selection)
                    self.lb.place(x=self.winfo_x(), y=self.winfo_y() + self.winfo_height(), width=self.winfo_width())
                    self.lb_up = True
                self.lb.delete(0, END)
                for w in words:
                    self.lb.insert(END, w)
                self.lb.config(height=len(words))  # 리스트박스의 높이를 항목 수에 맞게 조정
            else:
                if self.lb_up:
                    self.lb.destroy()
                    self.lb_up = False
        self.icursor(END)  # 커서를 항상 텍스트 끝으로 이동

    def selection(self, event):
        if self.lb_up:
            self.var.set(self.lb.get(self.lb.curselection()))
            self.lb.destroy()
            self.lb_up = False
            self.icursor(END)
            update_courses(self.var.get())

    def move_up(self, event):
        if self.lb_up:
            if self.lb.curselection() == ():
                index = '0'
            else:
                index = self.lb.curselection()[0]
            if index != '0':
                self.lb.selection_clear(first=index)
                index = str(int(index) - 1)
                self.lb.selection_set(first=index)
                self.lb.activate(index)

    def move_down(self, event):
        if self.lb_up:
            if self.lb.curselection() == ():
                index = '0'
            else:
                index = self.lb.curselection()[0]
            if index != END:
                self.lb.selection_clear(first=index)
                index = str(int(index) + 1)
                self.lb.selection_set(first=index)
                self.lb.activate(index)

    def comparison(self):
        pattern = self.var.get().lower()
        return [w for w in self.autocomplete_list if pattern in w.lower()]

def center_window(window):
    window.update_idletasks()
    width = 640
    height = 400
    x = (window.winfo_screenwidth() // 2) - (width // 2)
    y = (window.winfo_screenheight() // 2) - (height // 2)
    window.geometry(f'{width}x{height}+{x}+{y}')

def select_file(entry_widget, file_type):
    if file_type == "txt":
        file_path = filedialog.askopenfilename(filetypes=[("Text files", "*.txt")])
    elif file_type == "ppt":
        file_path = filedialog.askopenfilename(filetypes=[("PowerPoint files", "*.pptx")])
    elif file_type == "xlsx":
        file_path = filedialog.askopenfilename(filetypes=[("Excel files", "*.xlsx")])
    elif file_type == "image":
        file_path = filedialog.askopenfilename(filetypes=[("Image files", "*.jpg;*.jpeg;*.png")])
    if file_path:
        entry_widget.delete(0, END)
        entry_widget.insert(0, file_path)

# API 호출 비용 계산 함수
def calculate_cost(response):
    cost_per_token = 0.00006  # 예시: 토큰당 0.00006 달러
    total_tokens = response['usage']['total_tokens']
    return total_tokens * cost_per_token

# 파일 처리 함수
def process_selected_file():
    input_file_path = input_entry_txt.get()
    excel_file_path = input_entry_xlsx.get()

    if not input_file_path or not excel_file_path:
        messagebox.showerror("Error", "모든 파일 경로를 입력하세요.")
        return
    
    if not os.path.isfile(input_file_path):
        messagebox.showerror("Error", "유효한 텍스트 파일 경로를 입력하세요.")
        return
    
    if not os.path.isfile(excel_file_path):
        messagebox.showerror("Error", "유효한 Excel 파일 경로를 입력하세요.")
        return

    input_dir, input_filename = os.path.split(input_file_path)
    output_filename = f"re_{input_filename}"
    output_file_path = os.path.join(input_dir, output_filename)
    output_diff_file_path = os.path.join(input_dir, f"diff_{input_filename}")

    # 첫 번째 파일 읽기
    try:
        with open(input_file_path, 'r', encoding='utf-8') as file:
            text = file.read()
    except FileNotFoundError:
        print(f"파일을 찾을 수 없습니다: {input_file_path}")
        exit(1)

    # 텍스트를 분할하는 함수 (예: 단락 단위로 분할)
    def split_text(text, max_tokens=4000):
        words = text.split()
        chunks = []
        current_chunk = []
        current_length = 0

        for word in words:
            word_length = len(word)
            if current_length + word_length > max_tokens:
                chunks.append(' '.join(current_chunk) + ' //')  # 청크 끝에 '//' 추가
                current_chunk = [word]
                current_length = word_length
            else:
                current_chunk.append(word)
                current_length += word_length + 1  # +1 for space

        if current_chunk:
            chunks.append(' '.join(current_chunk) + ' //')  # 마지막 청크 끝에 '//' 추가

        return chunks

    # 텍스트를 분할
    text_chunks = split_text(text)

    # 분할된 텍스트를 처리하고 결과를 저장할 리스트
    processed_chunks = []

    # 각 청크를 처리
    for i, chunk in enumerate(text_chunks):
        print(f"Processing chunk {i+1}/{len(text_chunks)}")
        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "user", "content": f"다음 텍스트를 문단으로 나누고 줄바꿈을 추가해줘. 텍스트 내용을 변경하지않기, 내용 줄이지 않기, 소제목을 넣지 말기, 원본 내용을 유지해.: {chunk}"}
            ]
        )
        processed_chunk = response['choices'][0]['message']['content']
        processed_chunks.append(processed_chunk)

        # API 호출 비용 계산
        cost = calculate_cost(response)

    # 모든 처리된 청크를 결합
    final_text = '\n\n'.join(processed_chunks)

    # 결과를 같은 파일에 저장
    with open(output_file_path, 'w', encoding='utf-8') as file:
        file.write(final_text)

    # 두 번째 파일 읽기 및 수정된 파일명
    pptx_file_path = pptx_entry.get()
    image_path = image_entry.get()
    input_dir, input_filename = os.path.split(pptx_file_path)
    output_pptx_filename = f"updated_{input_filename}"
    output_pptx_path = os.path.join(input_dir, output_pptx_filename)

    try:
        with open(output_file_path, 'r', encoding='utf-8') as file:
            text = file.read()
    except FileNotFoundError:
        messagebox.showerror("Error", f"파일을 찾을 수 없습니다: {output_file_path}")
        return

    # 문단을 리스트로 분할
    paragraphs = text.split('\n\n')

    # OpenAI를 사용하여 소제목과 키워드 생성
    def generate_subtitle(paragraph):
        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "user", "content": f"다음 문단에 대한 소제목을 따옴표없이 한 개만 만들어 주세요: {paragraph}"}
            ]
        )
        return response['choices'][0]['message']['content'].strip()

    def generate_keywords(paragraph

):
        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "user", "content": f"다음 문단에 대한 키워드 블릿을 달고 최대 1~3개만 만들어 주세요: {paragraph}"}
            ]
        )
        return response['choices'][0]['message']['content'].strip().split(',')

    # 프레젠테이션 열기
    prs = Presentation(pptx_file_path)

    # 각 문단을 새로운 슬라이드로 추가
    for paragraph in paragraphs:
        # 슬라이드 레이아웃 설정 (1은 제목과 내용 슬라이드 레이아웃)
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # 소제목과 키워드 생성
        subtitle = generate_subtitle(paragraph)
        keywords = generate_keywords(paragraph)

        # 플레이스홀더에 텍스트 설정
        try:
            # 텍스트 설정
            slide.placeholders[22].text = subtitle  # 소제목
            slide.placeholders[17].text = ', '.join(keywords)  # 키워드
            slide.placeholders[16].text = paragraph.strip()  # 문단 내용

            # 이미지 추가
            placeholder = slide.placeholders[23]

            # 플레이스홀더의 위치와 크기 가져오기
            left = placeholder.left
            top = placeholder.top
            width = placeholder.width
            height = placeholder.height

            # 기존 플레이스홀더 삭제
            sp = placeholder._element
            sp.getparent().remove(sp)

            # 새 이미지 추가
            slide.shapes.add_picture(image_path, left, top, width, height)

        except (KeyError, IndexError) as e:
            print(f"플레이스홀더에 문제가 있습니다: {e}")
        except Exception as e:
            print(f"슬라이드에 내용을 추가하는 중 오류가 발생했습니다: {e}")

    # 프레젠테이션 저장
    prs.save(output_pptx_path)

    # 파일 읽기 함수
    def read_file(file_path):
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except FileNotFoundError:
            print(f"파일을 찾을 수 없습니다: {file_path}")
            exit(1)

    # 원본 텍스트와 결과 텍스트 읽기
    original_text = read_file(input_file_path)
    processed_text = read_file(output_file_path)

    # 줄바꿈 문자 제거 함수
    def remove_newlines(text):
        return text.replace('\n\n', '').replace('\r', '')

    # 줄바꿈 문자를 제거한 텍스트
    original_text_no_newlines = remove_newlines(original_text)
    processed_text_no_newlines = remove_newlines(processed_text)

    # 원본 텍스트와 결과 텍스트 비교 및 차이 출력
    def compare_texts_and_show_diff(original, processed, output_diff_file):
        if original == processed:
            message = "원본 텍스트와 결과 텍스트가 동일합니다."
        else:
            diff = list(difflib.ndiff(original, processed))
            added = [line[2:] for line in diff if line.startswith('+ ')]
            removed = [line[2:] for line in diff if line.startswith('- ')]
            
            added_message = "추가된 내용: " + ' '.join(added) if added else "추가된 내용이 없습니다."
            removed_message = "삭제된 내용: " + ' '.join(removed) if removed else "삭제된 내용이 없습니다."
            
            message = f"원본 텍스트와 결과 텍스트가 다릅니다. 차이점은 다음과 같습니다:\n\n{added_message}\n\n{removed_message}"
            
            # 차이점을 파일로 저장
            with open(output_diff_file, 'w', encoding='utf-8') as file:
                file.write(message)
        
        # 메시지 박스 표시
        # messagebox.showinfo("텍스트 비교 결과", message)

    # 비교 수행 및 차이점 출력
    compare_texts_and_show_diff(original_text_no_newlines, processed_text_no_newlines, output_diff_file_path)

    # process_files 함수 호출
    process_files(output_pptx_path, cost)

def process_files(pptx_file_path, api_cost):
    excel_file_path = input_entry_xlsx.get()
    save_pptx_path = pptx_file_path
    save_excel_path = excel_file_path
    
    institution_name = autocomplete_entry.get()
    course_name = course_name_combobox.get()
    week_number = week_number_entry.get()
    user_name = user_name_entry.get()

    execution_date = datetime.now().strftime("%Y-%m-%d")

    # 프레젠테이션 파일 열기
    prs = Presentation(pptx_file_path)

    # 1번 슬라이드 선택
    slide = prs.slides[0]

    # 기관명(대학교) - 기관명(대학교)에 값넣기
    title = slide.placeholders[14]  # 제목
    title.text = institution_name  # 사용자 입력 내용으로 변경

    # 과정명
    subtitle = slide.placeholders[0]  # 부제목
    subtitle.text = course_name  # 사용자 입력 내용으로 변경

    # 변경된 프레젠테이션 저장
    prs.save(save_pptx_path)

    # 기존 엑셀 파일 열기
    df = pd.read_excel(excel_file_path)

    # 새로운 데이터를 데이터프레임으로 만들기
    new_data = {
        "번호": [len(df) + 1],  # 번호는 현재 데이터프레임의 길이 + 1
        "기관명(학교명)": [institution_name],
        "과정명": [course_name],
        "차시(주차)": [week_number],
        "사용자": [user_name],
        "실행날짜": [execution_date],
        "최종결과물 위치": [save_pptx_path],
        "비용": [f"{api_cost:.2f}$"],
    }

    new_df = pd.DataFrame(new_data)

    # 기존 데이터프레임에 새로운 데이터 추가
    df = pd.concat([df, new_df], ignore_index=True)

    # 엑셀 파일로 저장
    df.to_excel(save_excel_path, index=False)

    # 엑셀 파일 열 너비 조정 (첫 번째 행 기준)
    wb = load_workbook(save_excel_path)
    ws = wb.active

    for col in ws.columns:
        max_length = 0
        column = col[0].column_letter  # Get the column name
        for cell in col:
            try:
                if len(str(cell.value)) > max_length:
                    max_length = len(cell.value)
            except:
                pass
        adjusted_width = (max_length + 2)
        ws.column_dimensions[column].width = adjusted_width

    wb.save(save_excel_path)

    print(f"변경된 프레젠테이션이 {save_pptx_path}에 저장되었습니다.")
    print(f"추가 정보가 엑셀 파일에 저장되었습니다.")
    messagebox.showinfo("완료", "모든 작업이 성공적으로 완료되었습니다.")

def update_courses(institution):
    courses = courses_df[courses_df['institution'] == institution]['course'].tolist()
    course_name_combobox['values'] = courses
    if courses:
        course_name_combobox.current(0)

# CSV 파일에서 데이터 읽기
df_institutions = pd.read_csv('institutions.csv')
institutions = df_institutions['institution'].tolist()

# 과정명 CSV 파일에서 데이터 읽기
courses_df = pd.read_csv('institutions_courses.csv')


# GUI 설정
root = tk.Tk()
root.title("문단 정리")

center_window(root)

frame = tk.Frame(root)
frame.pack(padx=10, pady=10, fill=tk.BOTH, expand=True)

font_large = ("Arial", 10)  # 기본 폰트 크기

# 텍스트 파일 선택
input_label_txt = tk.Label(frame, text="텍스트 파일 경로:", font=font_large)
input_label_txt.grid(row=0, column=0, padx=5, pady=5)

input_entry_txt = tk.Entry(frame, width=50, font=font_large)
input_entry_txt.grid(row=0, column=1, padx=5, pady=5)

select_button_txt = tk.Button(frame, text="파일 선택", command=lambda: select_file(input_entry_txt, "txt"), font=font_large)
select_button_txt.grid(row=0, column=2, padx=5, pady=5)

# PPT 파일 선택
input_label_ppt = tk.Label(frame, text="PowerPoint 파일 경로:", font=font_large)
input_label_ppt.grid(row=1, column=0, padx=5, pady=5)

pptx_entry = tk.Entry(frame, width=50, font=font_large)
pptx_entry.grid(row=1, column=1, padx=5, pady=5)

select_button_ppt = tk.Button(frame, text="파일 선택", command=lambda: select_file(pptx_entry, "ppt"), font=font_large)
select_button_ppt.grid(row=1, column=2, padx=5, pady=5)

# 이미지 파일 선택
input_label_image = tk.Label(frame, text="이미지 파일 경로:", font=font_large)
input_label_image.grid(row=2, column=0, padx=5, pady=5)

image_entry = tk.Entry(frame, width=50, font=font_large)
image_entry.grid(row=2, column=1, padx=5, pady=5)

select_button_image = tk.Button(frame, text="파일 선택", command=lambda: select_file(image_entry, "image"), font=font_large)
select_button_image.grid(row=2, column=2, padx=5, pady=5)

# Excel 파일 선택
input_label_xlsx = tk.Label(frame, text="Excel 파일 경로:", font=font_large)
input_label_xlsx.grid(row=3, column=0, padx=5, pady=5)

input_entry_xlsx = tk.Entry(frame, width=50, font=font_large)
input_entry_xlsx.grid(row=3, column=1, padx=5, pady=5)

select_button_xlsx = tk.Button(frame, text="파일 선택", command=lambda: select_file(input_entry_xlsx, "xlsx"), font=font_large)
select_button_xlsx.grid(row=3, column=2, padx=5, pady=5)

# 기관명 입력
tk.Label(frame, text="기관명(학교명):", font=font_large).grid(row=4, column=0, padx=5, pady=5)
autocomplete_entry = AutocompleteEntry(institutions, frame, width=50, font=font_large)
autocomplete_entry.grid(row=4, column=1, padx=5, pady=5)

# 과정명 입력
tk.Label(frame, text="과정명:", font=font_large).grid(row=5, column=0, padx=5, pady=5)
course_name_combobox = ttk.Combobox(frame, width=47, font=font_large)
course_name_combobox.grid(row=5, column=1, padx=5, pady=5)

# 차시(주차) 입력
tk.Label(frame, text="차시(주차):", font=font_large).grid(row=6, column=0, padx=5, pady=5)
week_number_entry = tk.Entry(frame, width=50, font=font_large)
week_number_entry.grid(row=6, column=1, padx=5, pady=5)

# 사용자 입력
tk.Label(frame, text="사용자:", font=font_large).grid(row=7, column=0, padx=5, pady=5)
user_name_entry = tk.Entry(frame, width=50, font=font_large)
user_name_entry.grid(row=7, column=1, padx=5, pady=5)

# 처리 버튼 가운데 정렬
process_button = tk.Button(frame, width=10, text="처리", command=process_selected_file, font=font_large)
process_button.grid(row=8, column=0, columnspan=3, pady=(20, 10))

root.mainloop()
