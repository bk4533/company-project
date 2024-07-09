import tkinter as tk
from tkinter import filedialog, messagebox
import os
import difflib
import openai
from pptx import Presentation
from dotenv import load_dotenv

# .env 파일의 환경 변수를 로드합니다.
load_dotenv()

# API 키를 환경 변수에서 가져옵니다.
API_KEY = os.getenv("FLASK_API_KEY")

# OpenAI API 키 설정
openai.api_key = API_KEY

def center_window(window):
    window.update_idletasks()
    width = 640
    height = 200
    x = (window.winfo_screenwidth() // 2) - (width // 2)
    y = (window.winfo_screenheight() // 2) - (height // 2)
    window.geometry(f'{width}x{height}+{x}+{y}')

# 파일 선택 함수
def select_file(entry_widget, file_type):
    if file_type == "txt":
        file_path = filedialog.askopenfilename(filetypes=[("Text files", "*.txt")])
    elif file_type == "ppt":
        file_path = filedialog.askopenfilename(filetypes=[("PowerPoint files", "*.pptx")])
    if file_path:
        entry_widget.delete(0, tk.END)
        entry_widget.insert(0, file_path)

# 파일 처리 함수
def process_selected_file():
    input_file_path = input_entry_txt.get()
    pptx_file_path = input_entry_ppt.get()

    if not input_file_path or not pptx_file_path:
        messagebox.showerror("Error", "모든 파일 경로를 입력하세요.")
        return
    
    if not os.path.isfile(input_file_path):
        messagebox.showerror("Error", "유효한 텍스트 파일 경로를 입력하세요.")
        return

    if not os.path.isfile(pptx_file_path):
        messagebox.showerror("Error", "유효한 PowerPoint 파일 경로를 입력하세요.")
        return

    input_dir, input_filename = os.path.split(input_file_path)
    output_filename = f"re_{input_filename}"
    output_file_path = os.path.join(input_dir, output_filename)
    output_diff_file_path = os.path.join(input_dir, f"diff_{input_filename}")

    # 첫 번째 파일 읽기
    try:
        with open(input_file_path, 'r', encoding='utf-8') as file:
            text = file.read()
    except FileNotFoundError:
        print(f"파일을 찾을 수 없습니다: {input_file_path}")
        exit(1)

    # 텍스트를 분할하는 함수 (예: 단락 단위로 분할)
    def split_text(text, max_tokens=4000):
        words = text.split()
        chunks = []
        current_chunk = []
        current_length = 0

        for word in words:
            word_length = len(word)
            if current_length + word_length > max_tokens:
                chunks.append(' '.join(current_chunk))
                current_chunk = [word]
                current_length = word_length
            else:
                current_chunk.append(word)
                current_length += word_length + 1  # +1 for space

        if current_chunk:
            chunks.append(' '.join(current_chunk))
        
        return chunks

    # 텍스트를 분할
    text_chunks = split_text(text)

    # 분할된 텍스트를 처리하고 결과를 저장할 리스트
    processed_chunks = []

    # 각 청크를 처리
    for i, chunk in enumerate(text_chunks):
        print(f"Processing chunk {i+1}/{len(text_chunks)}")
        response = openai.ChatCompletion.create(
            model="gpt-4o",
            messages=[
                {"role": "user", "content": f"다음 텍스트를 문단으로 나누고 줄바꿈을 추가해줘. 텍스트 내용을 변경하지않기, 내용 줄이지 않기, 소제목을 넣지 말기, 원본 내용을 유지해.: {chunk}"}
            ]
        )
        processed_chunk = response['choices'][0]['message']['content']
        processed_chunks.append(processed_chunk)

    # 모든 처리된 청크를 결합
    final_text = '\n\n'.join(processed_chunks)

    # 결과를 같은 파일에 저장
    with open(output_file_path, 'w', encoding='utf-8') as file:
        file.write(final_text)

    messagebox.showinfo("Success", f"파일이 성공적으로 저장되었습니다: {output_file_path}")

    # 두 번째 파일 읽기
    try:
        with open(output_file_path, 'r', encoding='utf-8') as file:
            text = file.read()
    except FileNotFoundError:
        messagebox.showerror("Error", f"파일을 찾을 수 없습니다: {output_file_path}")
        return

    # 문단을 리스트로 분할
    paragraphs = text.split('\n\n')

    # OpenAI를 사용하여 소제목과 키워드 생성
    def generate_subtitle(paragraph):
        response = openai.ChatCompletion.create(
            model="gpt-4o",
            messages=[
                {"role": "user", "content": f"다음 문단에 대한 소제목을 따옴표없이 한 개만 만들어 주세요: {paragraph}"}
            ]
        )
        return response['choices'][0]['message']['content'].strip()

    def generate_keywords(paragraph):
        response = openai.ChatCompletion.create(
            model="gpt-4o",
            messages=[
                {"role": "user", "content": f"다음 문단에 대한 키워드 블릿을 달고 최대 1~3개만 만들어 주세요: {paragraph}"}
            ]
        )
        return response['choices'][0]['message']['content'].strip().split(',')

    # 프레젠테이션 열기
    prs = Presentation(pptx_file_path)

    # 각 문단을 새로운 슬라이드로 추가
    for paragraph in paragraphs:
        # 슬라이드 레이아웃 설정 (1은 제목과 내용 슬라이드 레이아웃)
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # 소제목과 키워드 생성
        subtitle = generate_subtitle(paragraph)
        keywords = generate_keywords(paragraph)

        # 플레이스홀더에 텍스트 설정
        try:
            slide.placeholders[22].text = subtitle  # 소제목
            slide.placeholders[17].text = ', '.join(keywords)  # 키워드
            slide.placeholders[16].text = paragraph.strip()  # 문단 내용
        except KeyError as e:
            print(f"플레이스홀더가 슬라이드에 없습니다: {e}")

    # 프레젠테이션 저장
    output_pptx_path = os.path.join(input_dir, f"updated_{os.path.basename(pptx_file_path)}")
    prs.save(output_pptx_path)

    messagebox.showinfo("Success", f"PowerPoint 파일이 성공적으로 저장되었습니다: {output_pptx_path}")

    # 파일 읽기 함수
    def read_file(file_path):
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except FileNotFoundError:
            print(f"파일을 찾을 수 없습니다: {file_path}")
            exit(1)

    # 원본 텍스트와 결과 텍스트 읽기
    original_text = read_file(input_file_path)
    processed_text = read_file(output_file_path)

    # 줄바꿈 문자 제거 함수
    def remove_newlines(text):
        return text.replace('\n\n', '').replace('\r', '')

    # 줄바꿈 문자를 제거한 텍스트
    original_text_no_newlines = remove_newlines(original_text)
    processed_text_no_newlines = remove_newlines(processed_text)

    # 원본 텍스트와 결과 텍스트 비교 및 차이 출력
    def compare_texts_and_show_diff(original, processed, output_diff_file):
        if original == processed:
            message = "원본 텍스트와 결과 텍스트가 동일합니다."
        else:
            diff = list(difflib.ndiff(original, processed))
            added = [line[2:] for line in diff if line.startswith('+ ')]
            removed = [line[2:] for line in diff if line.startswith('- ')]
            
            added_message = "추가된 내용: " + ' '.join(added) if added else "추가된 내용이 없습니다."
            removed_message = "삭제된 내용: " + ' '.join(removed) if removed else "삭제된 내용이 없습니다."
            
            message = f"원본 텍스트와 결과 텍스트가 다릅니다. 차이점은 다음과 같습니다:\n\n{added_message}\n\n{removed_message}"
            
            # 차이점을 파일로 저장
            with open(output_diff_file, 'w', encoding='utf-8') as file:
                file.write(message)
        
        # 메시지 박스 표시
        root = tk.Tk()
        root.withdraw()  # 숨김 상태의 Tkinter 루트 윈도우
        messagebox.showinfo("텍스트 비교 결과", message)
        root.destroy()

    # 비교 수행 및 차이점 출력
    compare_texts_and_show_diff(original_text_no_newlines, processed_text_no_newlines, output_diff_file_path)

# GUI 설정
root = tk.Tk()
root.title("문단 정리")

center_window(root)

frame = tk.Frame(root)
frame.pack(padx=10, pady=10, fill=tk.BOTH, expand=True)

font_large = ("Arial", 10)  # 기본 폰트 크기

# 텍스트 파일 선택
input_label_txt = tk.Label(frame, text="텍스트 파일 경로:", font=font_large)
input_label_txt.grid(row=1, column=0, padx=5, pady=(20, 10))

input_entry_txt = tk.Entry(frame, width=50, font=font_large)
input_entry_txt.grid(row=1, column=1, padx=5, pady=(20, 10))

select_button_txt = tk.Button(frame, text="파일 선택", command=lambda: select_file(input_entry_txt, "txt"), font=font_large)
select_button_txt.grid(row=1, column=2, padx=5, pady=(20, 10))

# PPT 파일 선택
input_label_ppt = tk.Label(frame, text="PowerPoint 파일 경로:", font=font_large)
input_label_ppt.grid(row=2, column=0, padx=5, pady=(20, 10))

input_entry_ppt = tk.Entry(frame, width=50, font=font_large)
input_entry_ppt.grid(row=2, column=1, padx=5, pady=(20, 10))

select_button_ppt = tk.Button(frame, text="파일 선택", command=lambda: select_file(input_entry_ppt, "ppt"), font=font_large)
select_button_ppt.grid(row=2, column=2, padx=5, pady=(20, 10))

# 처리 버튼 가운데 정렬
process_button = tk.Button(frame, width=10, text="처리", command=process_selected_file, font=font_large)
process_button.grid(row=3, column=0, columnspan=3, pady=(20, 10))

root.mainloop()
