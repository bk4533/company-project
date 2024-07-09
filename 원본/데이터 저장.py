import pandas as pd
from pptx import Presentation
from datetime import datetime
from openpyxl import load_workbook
from tkinter import Tk, Label, Entry, Button, filedialog, StringVar, Listbox, END, ttk
import re

class AutocompleteEntry(Entry):
    def __init__(self, autocomplete_list, *args, **kwargs):
        super().__init__(*args, **kwargs)
        self.autocomplete_list = autocomplete_list
        self.var = StringVar()
        self.config(textvariable=self.var)
        self.var.trace('w', self.changed)
        self.bind("<Right>", self.selection)
        self.bind("<Double-Button-1>", self.selection)
        self.bind("<Return>", self.selection)
        self.bind("<Up>", self.move_up)
        self.bind("<Down>", self.move_down)
        self.lb_up = False

    def changed(self, name, index, mode):
        if self.var.get() == '':
            if self.lb_up:
                self.lb.destroy()
                self.lb_up = False
        else:
            words = self.comparison()
            if words:
                if not self.lb_up:
                    self.lb = Listbox(self.master)
                    self.lb.bind("<Double-Button-1>", self.selection)
                    self.lb.bind("<Right>", self.selection)
                    self.lb.bind("<Return>", self.selection)
                    self.lb.place(x=self.winfo_x(), y=self.winfo_y() + self.winfo_height(), width=self.winfo_width())
                    self.lb_up = True
                self.lb.delete(0, END)
                for w in words:
                    self.lb.insert(END, w)
                self.lb.config(height=len(words))  # 리스트박스의 높이를 항목 수에 맞게 조정
            else:
                if self.lb_up:
                    self.lb.destroy()
                    self.lb_up = False

    def selection(self, event):
        if self.lb_up:
            self.var.set(self.lb.get(self.lb.curselection()))
            self.lb.destroy()
            self.lb_up = False
            self.icursor(END)
            update_courses(self.var.get())

    def move_up(self, event):
        if self.lb_up:
            if self.lb.curselection() == ():
                index = '0'
            else:
                index = self.lb.curselection()[0]
            if index != '0':
                self.lb.selection_clear(first=index)
                index = str(int(index) - 1)
                self.lb.selection_set(first=index)
                self.lb.activate(index)

    def move_down(self, event):
        if self.lb_up:
            if self.lb.curselection() == ():
                index = '0'
            else:
                index = self.lb.curselection()[0]
            if index != END:
                self.lb.selection_clear(first=index)
                index = str(int(index) + 1)
                self.lb.selection_set(first=index)
                self.lb.activate(index)

    def comparison(self):
        pattern = self.var.get().lower()
        return [w for w in self.autocomplete_list if pattern in w.lower()]

def center_window(window):
    window.update_idletasks()
    width = 640
    height = 250
    x = (window.winfo_screenwidth() // 2) - (width // 2)
    y = (window.winfo_screenheight() // 2) - (height // 2)
    window.geometry(f'{width}x{height}+{x}+{y}')

def select_pptx_file():
    file_path = filedialog.askopenfilename(filetypes=[("PowerPoint files", "*.pptx")])
    pptx_entry.delete(0, 'end')
    pptx_entry.insert(0, file_path)

def select_xlsx_file():
    file_path = filedialog.askopenfilename(filetypes=[("Excel files", "*.xlsx")])
    xlsx_entry.delete(0, 'end')
    xlsx_entry.insert(0, file_path)


def process_files():
    pptx_file_path = pptx_entry.get()
    excel_file_path = xlsx_entry.get()
    save_pptx_path = pptx_file_path
    save_excel_path = excel_file_path
    
    institution_name = autocomplete_entry.get()
    course_name = course_name_combobox.get()
    week_number = week_number_entry.get()
    user_name = user_name_entry.get()



    execution_date = datetime.now().strftime("%Y-%m-%d")

    # 프레젠테이션 파일 열기
    prs = Presentation(pptx_file_path)

    # 1번 슬라이드 선택
    slide = prs.slides[0]

    # 기관명(대학교) - 기관명(대학교)에 값넣기
    title = slide.placeholders[14] # 제목
    title.text = institution_name  # 사용자 입력 내용으로 변경

    # 과정명
    subtitle = slide.placeholders[0] # 부제목
    subtitle.text = course_name  # 사용자 입력 내용으로 변경

    # 변경된 프레젠테이션 저장
    prs.save(save_pptx_path)

    # 기존 엑셀 파일 열기
    df = pd.read_excel(excel_file_path)

    # 새로운 데이터를 데이터프레임으로 만들기
    new_data = {
        "기관명(학교명)": [institution_name],
        "과정명": [course_name],
        "차시(주차)": [week_number],
        "사용자": [user_name],
        "실행날짜": [execution_date],
        "최종결과물 위치": [save_pptx_path],
    }

    new_df = pd.DataFrame(new_data)

    # 기존 데이터프레임에 새로운 데이터 추가
    df = pd.concat([df, new_df], ignore_index=True)

    # 엑셀 파일로 저장
    df.to_excel(save_excel_path, index=False)

    # 엑셀 파일 열 너비 조정 (첫 번째 행 기준)
    wb = load_workbook(save_excel_path)
    ws = wb.active

    for col in ws.columns:
        max_length = 0
        column = col[0].column_letter # Get the column name
        for cell in col:
            try:
                if len(str(cell.value)) > max_length:
                    max_length = len(cell.value)
            except:
                pass
        adjusted_width = (max_length + 2)
        ws.column_dimensions[column].width = adjusted_width

    wb.save(save_excel_path)

    print(f"변경된 프레젠테이션이 {save_pptx_path}에 저장되었습니다.")
    print(f"추가 정보가 엑셀 파일에 저장되었습니다.")

def update_courses(institution):
    courses = courses_df[courses_df['institution'] == institution]['course'].tolist()
    course_name_combobox['values'] = courses
    if courses:
        course_name_combobox.current(0)

# CSV 파일에서 데이터 읽기
df_institutions = pd.read_csv('institutions.csv')
institutions = df_institutions['institution'].tolist()  # 수정된 부분

# 과정명 CSV 파일에서 데이터 읽기
courses_df = pd.read_csv('institutions_courses.csv')

# Tkinter 초기화
root = Tk()
root.title("파일 처리")

center_window(root)

Label(root, text="PowerPoint 파일 경로:").grid(row=0, column=0, padx=10, pady=5)
pptx_entry = Entry(root, width=50)
pptx_entry.grid(row=0, column=1, padx=10, pady=5)
Button(root, text="파일 선택", command=select_pptx_file).grid(row=0, column=2, padx=10, pady=5)

Label(root, text="Excel 파일 경로:").grid(row=1, column=0, padx=10, pady=5)
xlsx_entry = Entry(root, width=50)
xlsx_entry.grid(row=1, column=1, padx=10, pady=5)
Button(root, text="파일 선택", command=select_xlsx_file).grid(row=1, column=2, padx=10, pady=5)

Label(root, text="기관명(학교명):").grid(row=2, column=0, padx=10, pady=5)
autocomplete_entry = AutocompleteEntry(institutions, root, width=50)
autocomplete_entry.grid(row=2, column=1, padx=10, pady=5)

Label(root, text="과정명:").grid(row=3, column=0, padx=10, pady=5)
course_name_combobox = ttk.Combobox(root, width=47)
course_name_combobox.grid(row=3, column=1, padx=10, pady=5)

Label(root, text="차시(주차):").grid(row=4, column=0, padx=10, pady=5)
week_number_entry = Entry(root, width=50)
week_number_entry.grid(row=4, column=1, padx=10, pady=5)

Label(root, text="사용자:").grid(row=5, column=0, padx=10, pady=5)
user_name_entry = Entry(root, width=50)
user_name_entry.grid(row=5, column=1, padx=10, pady=5)

Button(root, text="처리", command=process_files).grid(row=7, column=1, padx=10, pady=20)

root.mainloop()
