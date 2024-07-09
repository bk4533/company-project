from pptx import Presentation

# 프레젠테이션 불러오기
presentation_path = "C:\\Users\\bk453\\Desktop\\자동화 프로젝트\\ㅂㅂㅈ\\chatgpt1_3\\새 폴더\\sample.pptx"
prs = Presentation(presentation_path)

# 슬라이드 선택 (슬라이드 2는 인덱스 1)
slide = prs.slides[1]

# 이미지 경로 설정
image_path = "C:\\Users\\bk453\\Desktop\\자동화 프로젝트\\ㅂㅂㅈ\\chatgpt1_3\\새 폴더\\2.jpg"

try:
    # 플레이스홀더 인덱스 23 찾기
    placeholder = slide.placeholders[23]

    # 플레이스홀더의 위치와 크기 가져오기
    left = placeholder.left
    top = placeholder.top
    width = placeholder.width
    height = placeholder.height

    # 기존 플레이스홀더 삭제
    sp = placeholder._element
    sp.getparent().remove(sp)

    # 새 이미지 추가
    slide.shapes.add_picture(image_path, left, top, width, height)

except IndexError:
    print("플레이스홀더 인덱스 23이 슬라이드에 없습니다.")
except Exception as e:
    print(f"이미지를 교체하는 중 오류가 발생했습니다: {e}")

# 프레젠테이션 파일 저장 경로 설정
ppt_save_path = "C:\\Users\\bk453\\Desktop\\자동화 프로젝트\\ㅂㅂㅈ\\chatgpt1_3\\새 폴더\\sample_updated.pptx"

# 프레젠테이션 파일 저장
prs.save(ppt_save_path)
